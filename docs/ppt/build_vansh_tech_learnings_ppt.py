"""Builds a professional PPTX summarizing the technical learnings from the Vansh app debugging journey.

Usage (from repo root):
  python docs/ppt/build_vansh_tech_learnings_ppt.py

Optional screenshots:
  Put images in docs/ppt/images/ with these preferred filenames (any subset works):
    - 01_tree_fixed.png
    - 02_audio_upload_fixed.png
    - 03_settings_simplified.png
    - 04_network_error_before.png

If images are missing, the deck renders nice placeholders.

Output:
  docs/ppt/Vansh-Tech-Learnings.pptx
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "docs" / "ppt"
IMG_DIR = OUT_DIR / "images"
OUT_FILE = OUT_DIR / "Vansh-Tech-Learnings.pptx"


# --- Theme (matches Vansh aesthetic: khadi + suvarna) ---
KHADI_BG = RGBColor(249, 246, 238)      # warm off-white
CARD_BG = RGBColor(255, 252, 245)      # slightly brighter
INK = RGBColor(32, 33, 36)             # near-black
MUTED = RGBColor(110, 107, 101)        # warm gray
SUVARNA = RGBColor(196, 151, 52)       # gold
SUVARNA_DARK = RGBColor(140, 104, 26)  # deep gold
ACCENT_BLUE = RGBColor(41, 98, 255)    # for links/tech highlights


@dataclass(frozen=True)
class SlideImage:
    filename: str
    caption: str


def _add_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_header(slide, title: str, subtitle: Optional[str] = None) -> None:
    # Top band
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(1.1),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = SUVARNA
    band.line.fill.background()

    # Title
    t = slide.shapes.add_textbox(Inches(0.7), Inches(0.22), Inches(12), Inches(0.5))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.7), Inches(0.70), Inches(12), Inches(0.35))
        sf = s.text_frame
        sf.clear()
        sp = sf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.size = Pt(14)
        sr.font.color.rgb = RGBColor(255, 255, 255)


def _add_card(slide, x, y, w, h) -> None:
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = RGBColor(235, 228, 212)
    card.line.width = Pt(1)


def _add_bullets(slide, x, y, w, h, bullets: Iterable[str]) -> None:
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = b
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = INK


def _add_kpi(slide, x, y, label: str, value: str, color: RGBColor) -> None:
    pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, Inches(3.6), Inches(0.95))
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.fill.background()

    # Value
    vb = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.1), Inches(3.2), Inches(0.45))
    vtf = vb.text_frame
    vtf.clear()
    p = vtf.paragraphs[0]
    r = p.add_run()
    r.text = value
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)

    # Label
    lb = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.52), Inches(3.2), Inches(0.35))
    ltf = lb.text_frame
    ltf.clear()
    p2 = ltf.paragraphs[0]
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(12)
    r2.font.color.rgb = RGBColor(255, 255, 255)


def _try_add_image(slide, img_path: Path, x, y, w, h, caption: str) -> None:
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), x, y, width=w, height=h)
    else:
        # Placeholder
        ph = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
        ph.fill.solid()
        ph.fill.fore_color.rgb = RGBColor(245, 239, 225)
        ph.line.color.rgb = RGBColor(224, 212, 190)
        tx = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.25), w - Inches(0.6), Inches(0.7))
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "Drop screenshot here:\n" + img_path.name
        run.font.size = Pt(14)
        run.font.color.rgb = MUTED

    cap = slide.shapes.add_textbox(x, y + h + Inches(0.1), w, Inches(0.35))
    ctf = cap.text_frame
    ctf.clear()
    p = ctf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = caption
    r.font.size = Pt(12)
    r.font.color.rgb = MUTED


def _set_widescreen(prs: Presentation) -> None:
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def build() -> None:
    prs = Presentation()
    _set_widescreen(prs)

    # Use blank layout for full control
    blank = prs.slide_layouts[6]

    # --- Slide 1: Title ---
    s = prs.slides.add_slide(blank)
    _add_bg(s, KHADI_BG)
    _add_header(s, "Vansh App – Technical Learnings", "Debugging journey: Vriksha tree, API correctness, caching, and uploads")

    _add_card(s, Inches(0.8), Inches(1.6), Inches(11.8), Inches(2.4))
    title = s.shapes.add_textbox(Inches(1.2), Inches(1.95), Inches(11), Inches(1.6))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "From ‘flat & wrong’ → correct, dynamic, production-ready"
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = INK

    meta = s.shapes.add_textbox(Inches(1.2), Inches(3.2), Inches(11), Inches(0.6))
    mtf = meta.text_frame
    mtf.clear()
    p2 = mtf.paragraphs[0]
    p2.text = "React Native (Expo) • TypeScript • Express/MySQL • Graph layout algorithms"
    p2.font.size = Pt(16)
    p2.font.color.rgb = MUTED

    # --- Slide 2: What I built (system overview) ---
    s = prs.slides.add_slide(blank)
    _add_bg(s, KHADI_BG)
    _add_header(s, "System Overview", "Key pieces involved in the issues")

    _add_card(s, Inches(0.8), Inches(1.5), Inches(6.0), Inches(5.6))
    s.shapes.add_textbox(Inches(1.2), Inches(1.75), Inches(5.2), Inches(0.4)).text_frame.text = "Frontend (Expo)"
    _add_bullets(
        s,
        Inches(1.2),
        Inches(2.15),
        Inches(5.2),
        Inches(4.8),
        [
            "Vriksha: renders family tree layout + connectors",
            "Katha: audio record → multipart upload",
            "Settings: simplified UX",
            "State: Zustand stores + hooks",
        ],
    )

    _add_card(s, Inches(7.1), Inches(1.5), Inches(5.5), Inches(5.6))
    s.shapes.add_textbox(Inches(7.5), Inches(1.75), Inches(4.8), Inches(0.4)).text_frame.text = "Backend (Express + MySQL)"
    _add_bullets(
        s,
        Inches(7.5),
        Inches(2.15),
        Inches(4.8),
        Inches(4.8),
        [
            "Members + Relationships endpoints",
            "Relationship direction reversal per requester",
            "Cache-control to avoid stale data (304)",
            "DB constraints (ENUM) & normalization",
        ],
    )

    # --- Slide 3: Symptoms / outputs ---
    s = prs.slides.add_slide(blank)
    _add_bg(s, KHADI_BG)
    _add_header(s, "Observed Outputs", "What the app showed during debugging")

    imgs = [
        SlideImage("01_tree_fixed.png", "Tree after fixes (dynamic leveling + spouse/child placement)"),
        SlideImage("04_network_error_before.png", "Example error state captured during testing"),
    ]
    _try_add_image(s, IMG_DIR / imgs[0].filename, Inches(0.9), Inches(1.6), Inches(6.2), Inches(4.9), imgs[0].caption)
    _try_add_image(s, IMG_DIR / imgs[1].filename, Inches(7.3), Inches(1.6), Inches(5.1), Inches(4.9), imgs[1].caption)

    note = s.shapes.add_textbox(Inches(0.9), Inches(6.7), Inches(12.5), Inches(0.5))
    ntf = note.text_frame
    ntf.clear()
    p = ntf.paragraphs[0]
    p.text = "Screenshots auto-embed from docs/ppt/images (placeholders shown if missing)."
    p.font.size = Pt(12)
    p.font.color.rgb = MUTED

    # --- Slide 4: Primary challenges ---
    s = prs.slides.add_slide(blank)
    _add_bg(s, KHADI_BG)
    _add_header(s, "Challenges Encountered", "What broke and why it mattered")

    _add_card(s, Inches(0.8), Inches(1.55), Inches(12.0), Inches(5.9))
    _add_bullets(
        s,
        Inches(1.2),
        Inches(1.95),
        Inches(11.2),
        Inches(5.4),
        [
            "Family tree rendered as a flat row (no hierarchy)",
            "Tree appeared upside down (descendants at top)",
            "No relationships persisted in DB → ‘Relationships found: 0’",
            "Endpoint inconsistency: /members vs /members/:id/relationships returned different semantics",
            "HTTP 304 caching caused stale relationship types on device",
            "Audio upload failing (‘Network request failed’)",
        ],
    )

    # --- Slide 5: Root causes (technical) ---
    s = prs.slides.add_slide(blank)
    _add_bg(s, KHADI_BG)
    _add_header(s, "Root Causes", "What the debugging proved")

    _add_kpi(s, Inches(0.9), Inches(1.65), "Data correctness", "Relationships weren’t saved", SUVARNA_DARK)
    _add_kpi(s, Inches(4.85), Inches(1.65), "Semantics", "Type reversal missing", SUVARNA_DARK)
    _add_kpi(s, Inches(8.8), Inches(1.65), "Caching", "304 returned stale graph", SUVARNA_DARK)

    _add_card(s, Inches(0.8), Inches(2.8), Inches(12.0), Inches(4.7))
    _add_bullets(
        s,
        Inches(1.2),
        Inches(3.15),
        Inches(11.2),
        Inches(4.1),
        [
            "DB ENUM rejected detailed types like 'wife'/'son' → insert error, orphan members created",
            "getRelationships endpoint returned raw DB relationship_type (not reversed per requester)",
            "Client cached responses; UI used old relationship direction/type",
            "Tree can only be ‘dynamic’ if relationships are persisted + normalized and the graph is built from them",
        ],
    )

    # --- Slide 6: Solutions implemented ---
    s = prs.slides.add_slide(blank)
    _add_bg(s, KHADI_BG)
    _add_header(s, "Solutions Implemented", "Fixes across frontend, backend, and DB")

    _add_card(s, Inches(0.8), Inches(1.55), Inches(12.0), Inches(3.2))
    _add_bullets(
        s,
        Inches(1.2),
        Inches(1.9),
        Inches(11.2),
        Inches(2.9),
        [
            "Graph layout: BFS generation leveling (parents gen-1, children gen+1, spouse same gen)",
            "Backend: consistent relationship reversal in both endpoints",
            "Backend: Cache-Control headers to prevent stale 304 behavior",
            "Backend: map detailed relationship types → base ENUM; store detailed in relationship_subtype",
            "Frontend: audio upload fixed by not forcing multipart Content-Type (boundary required)",
        ],
    )

    _try_add_image(
        s,
        IMG_DIR / "02_audio_upload_fixed.png",
        Inches(0.9),
        Inches(4.95),
        Inches(6.2),
        Inches(2.05),
        "Katha upload working after multipart header fix",
    )
    _try_add_image(
        s,
        IMG_DIR / "03_settings_simplified.png",
        Inches(7.3),
        Inches(4.95),
        Inches(5.1),
        Inches(2.05),
        "Settings simplified to essentials",
    )

    # --- Slide 7: Algorithm detail (Vriksha) ---
    s = prs.slides.add_slide(blank)
    _add_bg(s, KHADI_BG)
    _add_header(s, "Tree Layout Algorithm", "Normalized graph → deterministic positioning")

    _add_card(s, Inches(0.8), Inches(1.55), Inches(12.0), Inches(5.9))

    bullets = [
        "Input: members + relationships (normalized: parent/child/spouse/sibling)",
        "Build adjacency lists for each node (parents, children, spouses)",
        "BFS from root: assign generations (Y levels) and normalize so minGen = 0",
        "Create spouse ‘family units’ and compute subtree widths bottom-up",
        "Assign X positions top-down using subtree widths (avoids overlap)",
        "Render orthogonal connectors (married + parent-child)",
    ]
    _add_bullets(s, Inches(1.2), Inches(1.95), Inches(11.2), Inches(5.3), bullets)

    # --- Slide 8: Lessons learned ---
    s = prs.slides.add_slide(blank)
    _add_bg(s, KHADI_BG)
    _add_header(s, "Key Technical Learnings", "What this taught me (practical takeaways)")

    _add_card(s, Inches(0.8), Inches(1.55), Inches(12.0), Inches(5.9))
    _add_bullets(
        s,
        Inches(1.2),
        Inches(1.95),
        Inches(11.2),
        Inches(5.3),
        [
            "Data model first: UI correctness depends on persisted relationship truth",
            "Normalize at the boundary: store base types + detailed subtype, not free-form strings",
            "API symmetry matters: all endpoints must agree on semantics",
            "Disable/handle caching during rapid iteration (304 can sabotage debugging)",
            "Multipart uploads: never hardcode boundary Content-Type in React Native",
            "Add logs where state transforms (DB → API → store → layout)",
        ],
    )

    # --- Slide 9: Next hardening steps ---
    s = prs.slides.add_slide(blank)
    _add_bg(s, KHADI_BG)
    _add_header(s, "Next Improvements", "To make this production-hard")

    _add_card(s, Inches(0.8), Inches(1.55), Inches(12.0), Inches(5.9))
    _add_bullets(
        s,
        Inches(1.2),
        Inches(1.95),
        Inches(11.2),
        Inches(5.3),
        [
            "Validate relationship creation server-side (reject cycles if desired)",
            "Enforce uniqueness constraints per direction + subtype rules",
            "Add integration tests for relationship semantics + reversal",
            "Add optimistic UI updates after add-member + relationship POST",
            "Add retry/backoff + clearer upload error details (status/body logging)",
        ],
    )

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_FILE))


if __name__ == "__main__":
    build()
    print(f"Wrote: {OUT_FILE}")
