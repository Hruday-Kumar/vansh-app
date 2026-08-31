"""
Builds a professional POC Pitch Deck for the Vansh Family Heritage App.

Usage (from repo root):
  pip install python-pptx   # if not already installed
  python docs/ppt/build_vansh_pitch_deck.py

Optional screenshots:
  Put images in docs/ppt/images/ with these filenames (any subset works):
    - app_home.png          (Time River home screen)
    - app_tree.png          (Family tree visualization)
    - app_memories.png      (Smriti memories gallery)
    - app_katha.png         (Katha voice recording)
    - app_vasiyat.png       (Vasiyat wisdom vault)
    - app_traditions.png    (Parampara traditions)
    - app_sharing.png       (QR code / sharing)
    - logo.png              (App logo)

Output:
  docs/ppt/Vansh-POC-Pitch-Deck.pptx
"""

from __future__ import annotations

from pathlib import Path
from typing import Iterable, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "docs" / "ppt"
IMG_DIR = OUT_DIR / "images"
OUT_FILE = OUT_DIR / "Vansh-POC-Pitch-Deck.pptx"

# ─── Theme (Digital Sanskriti) ───────────────────────────────────
KHADI      = RGBColor(249, 246, 238)    # warm off-white bg
CARD_BG    = RGBColor(255, 252, 245)    # card bg
INK        = RGBColor(32, 33, 36)       # near-black text
MUTED      = RGBColor(110, 107, 101)    # warm gray
SUVARNA    = RGBColor(196, 151, 52)     # temple gold
SUVARNA_DK = RGBColor(140, 104, 26)     # deep gold
SINDOOR    = RGBColor(196, 83, 42)      # vermilion
PADMA      = RGBColor(196, 107, 128)    # lotus pink
NEELAM     = RGBColor(41, 98, 255)      # sapphire blue
PANNA      = RGBColor(34, 139, 94)      # emerald
HALDI      = RGBColor(218, 165, 32)     # turmeric
WHITE      = RGBColor(255, 255, 255)
DARK_BG    = RGBColor(42, 40, 38)       # dark slide bg
SLIDE_W    = Inches(13.333)
SLIDE_H    = Inches(7.5)


# ─── Helpers ─────────────────────────────────────────────────────
def _set_widescreen(prs: Presentation) -> None:
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H


def _add_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, x, y, w, h, fill_color, line_color=None, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _add_text(slide, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return tf


def _add_multi_text(slide, x, y, w, h, lines: list[tuple[str, int, RGBColor, bool]], align=PP_ALIGN.LEFT):
    """Add a textbox with multiple formatted lines."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    first = True
    for text, size, color, bold in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
    return tf


def _add_header_band(slide, title: str, subtitle: str = "", bg_color=SUVARNA):
    band = _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), bg_color, radius=False)
    _add_text(slide, Inches(0.7), Inches(0.18), Inches(12), Inches(0.5), title, 30, WHITE, True)
    if subtitle:
        _add_text(slide, Inches(0.7), Inches(0.68), Inches(12), Inches(0.35), subtitle, 14, WHITE)


def _add_slide_number(slide, num: int, total: int):
    _add_text(slide, Inches(12.2), Inches(7.1), Inches(1), Inches(0.3),
              f"{num}/{total}", 10, MUTED, align=PP_ALIGN.RIGHT)


def _add_bullets(slide, x, y, w, h, bullets: list[str], size=17, color=INK, spacing=Pt(8)):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = b
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = spacing


def _add_kpi_pill(slide, x, y, value: str, label: str, color: RGBColor, w=Inches(3.6)):
    _add_shape(slide, x, y, w, Inches(1.05), color)
    _add_text(slide, x + Inches(0.2), y + Inches(0.1), w - Inches(0.4), Inches(0.45),
              value, 24, WHITE, True, PP_ALIGN.CENTER)
    _add_text(slide, x + Inches(0.2), y + Inches(0.58), w - Inches(0.4), Inches(0.35),
              label, 12, WHITE, False, PP_ALIGN.CENTER)


def _add_icon_card(slide, x, y, w, h, icon: str, title: str, desc: str, accent: RGBColor):
    """Compact feature card with emoji icon."""
    _add_shape(slide, x, y, w, h, CARD_BG, RGBColor(235, 228, 212))
    # Accent strip on left
    _add_shape(slide, x, y, Inches(0.08), h, accent, radius=False)
    # Icon
    _add_text(slide, x + Inches(0.2), y + Inches(0.1), Inches(0.5), Inches(0.5), icon, 22, INK, align=PP_ALIGN.CENTER)
    # Title
    _add_text(slide, x + Inches(0.7), y + Inches(0.12), w - Inches(0.9), Inches(0.35), title, 14, INK, True)
    # Description
    _add_text(slide, x + Inches(0.7), y + Inches(0.45), w - Inches(0.9), h - Inches(0.55), desc, 11, MUTED)


def _try_add_image(slide, filename: str, x, y, w, h, caption: str = ""):
    img_path = IMG_DIR / filename
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), x, y, width=w, height=h)
    else:
        ph = _add_shape(slide, x, y, w, h, RGBColor(245, 239, 225), RGBColor(224, 212, 190))
        _add_text(slide, x + Inches(0.2), y + h / 2 - Inches(0.3), w - Inches(0.4), Inches(0.7),
                  f"[Screenshot: {filename}]", 12, MUTED, align=PP_ALIGN.CENTER)
    if caption:
        _add_text(slide, x, y + h + Inches(0.05), w, Inches(0.3), caption, 10, MUTED, align=PP_ALIGN.CENTER)


def _add_table_row_shapes(slide, x, y, w, cols: list[tuple[str, float]], row_h, bg, text_color, bold=False, size=13):
    """Draw table row using shapes (more control than pptx tables)."""
    cx = x
    for text, col_w_frac in cols:
        cw = w * col_w_frac
        _add_shape(slide, cx, y, cw, row_h, bg, RGBColor(220, 215, 205), radius=False)
        _add_text(slide, cx + Inches(0.1), y + Inches(0.03), cw - Inches(0.2), row_h - Inches(0.06),
                  text, size, text_color, bold, anchor=MSO_ANCHOR.MIDDLE)
        cx += cw


# ═════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═════════════════════════════════════════════════════════════════

TOTAL_SLIDES = 15


def slide_01_title(prs, blank):
    """Title slide."""
    s = prs.slides.add_slide(blank)
    _add_bg(s, DARK_BG)

    # Gold accent band top
    _add_shape(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), SUVARNA, radius=False)

    # Title block
    _add_multi_text(s, Inches(1.5), Inches(1.4), Inches(10.5), Inches(3.5), [
        ("🪷", 50, SUVARNA, False),
        ("VANSH  (वंश)", 54, WHITE, True),
        ("", 12, WHITE, False),
        ("Preserving Family Legacies Across Generations", 24, SUVARNA, False),
        ("", 10, WHITE, False),
        ("A mobile-first platform to document, share, and preserve", 18, RGBColor(200, 195, 185), False),
        ("your family's stories, memories, traditions, and wisdom forever.", 18, RGBColor(200, 195, 185), False),
    ], PP_ALIGN.CENTER)

    # Tech badges
    badges = "React Native  •  TypeScript  •  Node.js/Express  •  MySQL  •  Gemini AI  •  Firebase"
    _add_text(s, Inches(1.5), Inches(5.3), Inches(10.5), Inches(0.4), badges, 13, MUTED, align=PP_ALIGN.CENTER)

    # Bottom bar
    _add_shape(s, Inches(0), Inches(6.8), SLIDE_W, Inches(0.7), SUVARNA_DK, radius=False)
    _add_text(s, Inches(1), Inches(6.88), Inches(5), Inches(0.4), "POC Presentation  •  February 2026", 14, WHITE)
    _add_text(s, Inches(7), Inches(6.88), Inches(5.5), Inches(0.4), "Version 2.0  •  MIT License", 14, WHITE, align=PP_ALIGN.RIGHT)

    _add_slide_number(s, 1, TOTAL_SLIDES)


def slide_02_problem(prs, blank):
    """The Problem."""
    s = prs.slides.add_slide(blank)
    _add_bg(s, KHADI)
    _add_header_band(s, "The Problem", "Why families are losing their heritage", SINDOOR)

    cards = [
        ("📱", "Digital Disconnect",
         "Younger generations are losing connection with family elders. Nuclear family trend means stories, recipes, and traditions are not being passed down.",
         SINDOOR),
        ("💬", "Unstructured Communication",
         "Family memories scattered across WhatsApp groups, Google Photos, and phone galleries. No context, no relationships, content gets buried over time.",
         RGBColor(180, 80, 60)),
        ("🌍", "Diaspora Challenge",
         "32M+ Indians abroad losing cultural roots. No tool helps them preserve & share heritage in an Indian cultural context with proper kinship terminology.",
         RGBColor(160, 60, 40)),
    ]

    for i, (icon, title, desc, accent) in enumerate(cards):
        _add_icon_card(s, Inches(0.8), Inches(1.5 + i * 2.0), Inches(7.5), Inches(1.8), icon, title, desc, accent)

    # Stats on right
    _add_shape(s, Inches(8.8), Inches(1.5), Inches(4.0), Inches(5.8), CARD_BG, RGBColor(235, 228, 212))
    _add_text(s, Inches(9.0), Inches(1.7), Inches(3.6), Inches(0.4), "The Urgency", 18, SINDOOR, True, PP_ALIGN.CENTER)

    stats = [
        ("80%", "of oral traditions lost\nwithin 2 generations"),
        ("300M+", "Indian households with\nno heritage preservation tool"),
        ("72%", "of millennials wish they\nknew more family history"),
        ("0", "apps with Indian kinship\nterms + voice stories + tree"),
    ]
    for i, (val, desc) in enumerate(stats):
        _add_text(s, Inches(9.2), Inches(2.3 + i * 1.25), Inches(1.2), Inches(0.4), val, 28, SUVARNA_DK, True)
        _add_text(s, Inches(10.3), Inches(2.35 + i * 1.25), Inches(2.3), Inches(1.0), desc, 12, MUTED)

    _add_slide_number(s, 2, TOTAL_SLIDES)


def slide_03_solution(prs, blank):
    """Our Solution."""
    s = prs.slides.add_slide(blank)
    _add_bg(s, KHADI)
    _add_header_band(s, "Our Solution", "Vansh — A digital family heritage vault", PANNA)

    _add_text(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.7),
              "One app that combines Family Tree + Voice Stories + Photo Memories + Traditions + Time-Locked Messages — all designed for Indian families.",
              18, INK)

    # Feature grid 2x3
    features = [
        ("🌳", "Vriksha (Family Tree)", "Interactive tree with 40+ Hindi kinship terms, gesture navigation, relationship finder", PANNA),
        ("📸", "Smriti (Memories)", "Photo/video gallery with AI tagging, member tagging, mosaic layout, event albums", NEELAM),
        ("🎙️", "Katha (Stories)", "Voice recordings with waveform viz, transcription, photo-story stitching", SUVARNA_DK),
        ("🪔", "Parampara (Traditions)", "Recipes, rituals, festivals with step-by-step docs and cultural context", SINDOOR),
        ("💌", "Vasiyat (Wisdom Vault)", "End-to-end encrypted time-locked messages with 4 trigger types", PADMA),
        ("🎉", "Nimantran (Invitations)", "15 ceremony types, RSVP tracking, reminders for family events", HALDI),
    ]

    for i, (icon, title, desc, accent) in enumerate(features):
        col = i % 3
        row = i // 3
        x = Inches(0.8) + col * Inches(4.1)
        y = Inches(2.3) + row * Inches(2.2)
        _add_icon_card(s, x, y, Inches(3.9), Inches(2.0), icon, title, desc, accent)

    _add_slide_number(s, 3, TOTAL_SLIDES)


def slide_04_demo_home(prs, blank):
    """Product — Time River & Home."""
    s = prs.slides.add_slide(blank)
    _add_bg(s, KHADI)
    _add_header_band(s, "Product: Time River (Home Feed)", "Chronological feed of all family heritage content")

    _try_add_image(s, "app_home.png", Inches(0.8), Inches(1.5), Inches(4.5), Inches(5.5), "Time River — Family Feed")

    _add_shape(s, Inches(5.8), Inches(1.5), Inches(7.0), Inches(5.5), CARD_BG, RGBColor(235, 228, 212))
    _add_bullets(s, Inches(6.2), Inches(1.8), Inches(6.2), Inches(5.0), [
        "🏠  Unified chronological feed of ALL family content",
        "📸  Memories, voice stories, traditions — one timeline",
        "🔄  Cursor-based pagination for smooth scrolling",
        "🎵  Ambient audio mode — family stories play in background",
        "📊  Family statistics dashboard",
        "🔍  Quick search across all content types",
        "📱  Pull-to-refresh with haptic feedback",
        "🌐  Works offline — syncs when back online",
    ], 15, INK)

    _add_slide_number(s, 4, TOTAL_SLIDES)


def slide_05_demo_tree(prs, blank):
    """Product — Vriksha Family Tree."""
    s = prs.slides.add_slide(blank)
    _add_bg(s, KHADI)
    _add_header_band(s, "Product: Vriksha (Interactive Family Tree)", "Custom graph algorithm with Indian kinship system")

    _try_add_image(s, "app_tree.png", Inches(0.8), Inches(1.5), Inches(5.0), Inches(5.5), "Vriksha — Family Tree")

    _add_shape(s, Inches(6.3), Inches(1.5), Inches(6.5), Inches(5.5), CARD_BG, RGBColor(235, 228, 212))

    _add_text(s, Inches(6.6), Inches(1.7), Inches(6.0), Inches(0.4), "Technical Highlights", 16, SUVARNA_DK, True)

    _add_bullets(s, Inches(6.6), Inches(2.1), Inches(6.0), Inches(4.8), [
        "Custom 6-phase tree layout algorithm (BFS → positioning)",
        "40+ Hindi kinship terms (Sudanese kinship system)",
        "Auto-inference: 1 relationship → creates transitive links",
        "Pinch-to-zoom (0.3x–2.5x) + pan with silk-like inertia",
        "'View As' — re-root tree from any member's perspective",
        "Select 2 members → BFS finds & labels relationship",
        "Real-time Firebase sync across devices",
        "Closure table → O(1) ancestor/descendant queries",
        "Gesture: double-tap to center, haptic feedback",
    ], 13, INK, Pt(6))

    _add_slide_number(s, 5, TOTAL_SLIDES)


def slide_06_architecture(prs, blank):
    """System Architecture (text-based diagram)."""
    s = prs.slides.add_slide(blank)
    _add_bg(s, KHADI)
    _add_header_band(s, "System Architecture", "Client-Server with Offline-First capability")

    # ── Left: Mobile App box ──
    _add_shape(s, Inches(0.5), Inches(1.5), Inches(4.2), Inches(5.6), RGBColor(240, 248, 255), RGBColor(41, 98, 255))
    _add_text(s, Inches(0.7), Inches(1.6), Inches(3.8), Inches(0.4), "📱 Mobile App (Expo/React Native)", 14, NEELAM, True, PP_ALIGN.CENTER)

    mobile_items = [
        ("Expo Router", "File-based navigation"),
        ("Zustand (9 stores)", "State management"),
        ("NativeWind", "Tailwind CSS styling"),
        ("Reanimated 4", "60fps animations"),
        ("SQLite", "Offline database"),
        ("Secure Store", "Token encryption"),
        ("Firebase RT DB", "Real-time sync"),
    ]
    for i, (name, desc) in enumerate(mobile_items):
        _add_shape(s, Inches(0.7), Inches(2.15 + i * 0.65), Inches(3.8), Inches(0.55), CARD_BG, RGBColor(200, 210, 230))
        _add_text(s, Inches(0.85), Inches(2.2 + i * 0.65), Inches(1.6), Inches(0.25), name, 11, INK, True)
        _add_text(s, Inches(0.85), Inches(2.42 + i * 0.65), Inches(3.4), Inches(0.2), desc, 9, MUTED)

    # ── Center: Arrows ──
    _add_text(s, Inches(4.85), Inches(3.5), Inches(0.8), Inches(0.5), "REST\nAPI", 12, MUTED, align=PP_ALIGN.CENTER)
    _add_text(s, Inches(4.85), Inches(4.2), Inches(0.8), Inches(0.4), "⟷", 24, SUVARNA, True, PP_ALIGN.CENTER)
    _add_text(s, Inches(4.85), Inches(4.6), Inches(0.8), Inches(0.3), "JSON\nHTTPS", 10, MUTED, align=PP_ALIGN.CENTER)

    # ── Right top: Backend ──
    _add_shape(s, Inches(5.9), Inches(1.5), Inches(3.8), Inches(3.5), RGBColor(255, 245, 238), RGBColor(196, 83, 42))
    _add_text(s, Inches(6.1), Inches(1.6), Inches(3.4), Inches(0.4), "⚙️ Backend (Node.js/Express)", 14, SINDOOR, True, PP_ALIGN.CENTER)

    backend_items = ["JWT Auth + bcrypt", "Rate Limiting (4-tier)", "Helmet.js Security", "Multer File Uploads", "Winston + Sentry"]
    for i, item in enumerate(backend_items):
        _add_shape(s, Inches(6.1), Inches(2.15 + i * 0.52), Inches(3.4), Inches(0.43), CARD_BG, RGBColor(230, 210, 200))
        _add_text(s, Inches(6.25), Inches(2.2 + i * 0.52), Inches(3.1), Inches(0.35), item, 11, INK)

    # ── Right bottom: Database ──
    _add_shape(s, Inches(5.9), Inches(5.3), Inches(3.8), Inches(1.8), RGBColor(240, 255, 240), RGBColor(34, 139, 94))
    _add_text(s, Inches(6.1), Inches(5.4), Inches(3.4), Inches(0.4), "🗄️ MySQL 8.0 Database", 14, PANNA, True, PP_ALIGN.CENTER)
    _add_bullets(s, Inches(6.1), Inches(5.8), Inches(3.4), Inches(1.2), [
        "13 tables • UUID PKs • JSON columns",
        "Closure table for O(1) ancestry",
        "19 indexes • CASCADE deletes",
    ], 11, INK, Pt(3))

    # ── Far right: AI & External ──
    _add_shape(s, Inches(10.1), Inches(1.5), Inches(2.8), Inches(2.5), RGBColor(255, 248, 240), RGBColor(196, 151, 52))
    _add_text(s, Inches(10.3), Inches(1.6), Inches(2.4), Inches(0.4), "🤖 AI (Gemini 1.5)", 14, SUVARNA_DK, True, PP_ALIGN.CENTER)
    _add_bullets(s, Inches(10.3), Inches(2.0), Inches(2.4), Inches(1.8), [
        "Image analysis",
        "Era detection",
        "Transcription",
        "Digital Echo",
    ], 11, INK, Pt(2))

    _add_shape(s, Inches(10.1), Inches(4.3), Inches(2.8), Inches(1.5), RGBColor(255, 245, 245), RGBColor(200, 100, 100))
    _add_text(s, Inches(10.3), Inches(4.4), Inches(2.4), Inches(0.4), "🔥 Firebase", 14, SINDOOR, True, PP_ALIGN.CENTER)
    _add_bullets(s, Inches(10.3), Inches(4.8), Inches(2.4), Inches(0.8), [
        "Realtime DB sync",
        "Cross-device collab",
    ], 11, INK, Pt(2))

    _add_shape(s, Inches(10.1), Inches(6.1), Inches(2.8), Inches(1.0), RGBColor(245, 240, 255), RGBColor(100, 80, 180))
    _add_text(s, Inches(10.3), Inches(6.2), Inches(2.4), Inches(0.4), "📁 File Storage", 13, RGBColor(100, 80, 180), True, PP_ALIGN.CENTER)
    _add_text(s, Inches(10.3), Inches(6.55), Inches(2.4), Inches(0.3), "Server disk /uploads/", 10, MUTED, align=PP_ALIGN.CENTER)

    _add_slide_number(s, 6, TOTAL_SLIDES)


def slide_07_tech_deep_dive(prs, blank):
    """Tech Stack Deep Dive."""
    s = prs.slides.add_slide(blank)
    _add_bg(s, KHADI)
    _add_header_band(s, "Technology Stack", "Production-grade choices across the stack")

    # Table header
    headers = [("Layer", 0.18), ("Technology", 0.32), ("Why This Choice", 0.50)]
    _add_table_row_shapes(s, Inches(0.5), Inches(1.4), Inches(12.3), headers, Inches(0.4), SUVARNA_DK, WHITE, True, 13)

    rows = [
        ("Frontend", "React Native 0.81 + Expo SDK 54", "Cross-platform (iOS/Android/Web), OTA updates, 70+ Expo modules"),
        ("Language", "TypeScript 5.9", "Type safety across frontend + backend, better DX, fewer runtime errors"),
        ("Navigation", "Expo Router (file-based)", "Next.js-style routing, deep linking built-in, type-safe"),
        ("Styling", "NativeWind (Tailwind CSS)", "Utility-first, familiar for web devs, compiled to StyleSheet"),
        ("Animations", "Reanimated 4 + Gesture Handler", "60fps UI-thread animations, gesture system, haptic feedback"),
        ("State", "Zustand 5 (9 stores)", "Lightweight (~1KB), no boilerplate, persist middleware"),
        ("Backend", "Node.js + Express 4", "Fast prototyping, massive ecosystem, TypeScript support"),
        ("Database", "MySQL 8.0", "ACID compliance, JSON columns, proven at scale"),
        ("Offline DB", "expo-sqlite", "Local SQLite for offline-first, sync queue"),
        ("AI/ML", "Google Gemini 1.5 Flash", "Image analysis, transcription, persona chat, low cost"),
        ("Auth", "JWT + bcrypt + Refresh Tokens", "Stateless auth, secure password hashing, session management"),
        ("Security", "Helmet + Rate Limiting + AES-256", "Production-grade headers, DDoS protection, E2E encryption"),
        ("Monitoring", "Sentry + Winston", "Error tracking, structured logging, PII filtering"),
        ("Real-time", "Firebase Realtime DB", "Cross-device sync, debounced updates, conflict resolution"),
        ("Testing", "Jest 30 + Testing Library", "Unit + component + API tests, pre-commit hooks"),
    ]

    for i, (layer, tech, reason) in enumerate(rows):
        bg = CARD_BG if i % 2 == 0 else KHADI
        cols = [(layer, 0.18), (tech, 0.32), (reason, 0.50)]
        _add_table_row_shapes(s, Inches(0.5), Inches(1.8 + i * 0.36), Inches(12.3), cols, Inches(0.36), bg, INK, False, 11)

    _add_slide_number(s, 7, TOTAL_SLIDES)


def slide_08_algorithms(prs, blank):
    """Key Algorithms & Innovations."""
    s = prs.slides.add_slide(blank)
    _add_bg(s, KHADI)
    _add_header_band(s, "Key Algorithms & Innovations", "What makes Vansh technically unique")

    # Left column
    _add_shape(s, Inches(0.5), Inches(1.4), Inches(6.2), Inches(5.8), CARD_BG, RGBColor(235, 228, 212))
    _add_text(s, Inches(0.7), Inches(1.5), Inches(5.8), Inches(0.4), "🌳  Family Tree Layout (6-Phase Pipeline)", 15, SUVARNA_DK, True)

    phases = [
        "1. BFS Generation Assignment — Traverse from root, assign Y-levels",
        "2. Family Unit Pairing — Group spouses, collect their children",
        "3. Unit Graph Construction — Directed parent→child unit edges",
        "4. Bottom-Up Subtree Width — DFS with cycle detection",
        "5. Top-Down Position Assignment — Center children under parents",
        "6. Connector Generation — Orthogonal bracket-style lines",
    ]
    _add_bullets(s, Inches(0.7), Inches(1.95), Inches(5.8), Inches(2.5), phases, 12, INK, Pt(4))

    _add_text(s, Inches(0.7), Inches(4.0), Inches(5.8), Inches(0.4), "🔍  Relationship Resolution (BFS + Pattern Matching)", 15, SUVARNA_DK, True)
    _add_bullets(s, Inches(0.7), Inches(4.45), Inches(5.8), Inches(2.5), [
        "BFS traversal up to depth 12 through adjacency list",
        "Edge sequences → 40+ Hindi kinship terms",
        "e.g. parent→spouse→parent = Sasur (ससुर, father-in-law)",
        "Generic fallback: '5x Great-Grandfather' for deep ancestry",
    ], 12, INK, Pt(4))

    # Right column
    _add_shape(s, Inches(7.0), Inches(1.4), Inches(5.8), Inches(2.7), CARD_BG, RGBColor(235, 228, 212))
    _add_text(s, Inches(7.2), Inches(1.5), Inches(5.4), Inches(0.4), "⚡  Auto-Inference Engine", 15, PANNA, True)
    _add_bullets(s, Inches(7.2), Inches(1.95), Inches(5.4), Inches(2.0), [
        "Adding spouse → auto-links as parent to existing children",
        "Adding child → auto-creates sibling + co-parent spouse",
        "Reduces manual input by ~60%",
        "Transitive closure maintained automatically",
    ], 12, INK, Pt(4))

    _add_shape(s, Inches(7.0), Inches(4.3), Inches(5.8), Inches(2.9), CARD_BG, RGBColor(235, 228, 212))
    _add_text(s, Inches(7.2), Inches(4.4), Inches(5.4), Inches(0.4), "💾  Offline-First Architecture", 15, NEELAM, True)
    _add_bullets(s, Inches(7.2), Inches(4.85), Inches(5.4), Inches(2.2), [
        "SQLite sync queue: pending → processing → completed/failed",
        "executeOrQueue: tries online first, queues if no network",
        "Network polling every 5 sec, auto-sync on reconnect",
        "LRU cache with stale-while-revalidate (same as Vercel)",
        "5 cache namespaces with independent TTLs",
    ], 12, INK, Pt(4))

    _add_slide_number(s, 8, TOTAL_SLIDES)


def slide_09_security(prs, blank):
    """Security & Privacy."""
    s = prs.slides.add_slide(blank)
    _add_bg(s, KHADI)
    _add_header_band(s, "Security & Privacy Architecture", "Family data demands bank-level protection", SINDOOR)

    layers = [
        ("🔐", "Authentication", "JWT access tokens + 30-day refresh tokens in DB. bcrypt 10-round password hashing.", SINDOOR),
        ("🛡️", "Transport Security", "Helmet.js headers (CSP, HSTS, X-Frame-Options). CORS whitelist in production.", RGBColor(180, 70, 50)),
        ("⚡", "Rate Limiting", "4-tier protection: General (100/15min), Auth (5/15min), Sensitive (10/hr), Upload (20/hr).", RGBColor(160, 60, 40)),
        ("🔒", "Data Encryption", "AES-256-GCM for family data. Separate encryption key for Vasiyat messages. expo-secure-store for keys.", RGBColor(140, 50, 30)),
        ("👥", "Multi-Tenancy", "Every DB query filtered by family_id. Users can NEVER access another family's data.", RGBColor(120, 40, 25)),
        ("🎭", "RBAC", "4 roles: Admin > Elder > Member > Viewer. Endpoint-level permission checks.", RGBColor(100, 35, 20)),
        ("📊", "Monitoring", "Sentry error tracking with PII stripping. Winston structured logging with file rotation.", RGBColor(85, 30, 18)),
    ]

    for i, (icon, title, desc, accent) in enumerate(layers):
        y = Inches(1.4 + i * 0.85)
        _add_icon_card(s, Inches(0.5), y, Inches(12.3), Inches(0.78), icon, title, desc, accent)

    _add_slide_number(s, 9, TOTAL_SLIDES)


def slide_10_ai_features(prs, blank):
    """AI-Powered Features."""
    s = prs.slides.add_slide(blank)
    _add_bg(s, KHADI)
    _add_header_band(s, "AI-Powered Heritage (Gemini 1.5 Flash)", "Making family heritage interactive and intelligent", SUVARNA)

    features = [
        ("📸", "Image Analysis",
         "Upload photo → Gemini extracts objects, people count, description, auto-tags. Enables smart search across decades of family photos.",
         SUVARNA_DK, Inches(0.5), Inches(1.5), Inches(6.0), Inches(1.6)),
        ("📅", "Era Detection",
         "AI estimates photo decade from clothing styles, technology visible, and photo quality — auto-places in family timeline.",
         HALDI, Inches(6.8), Inches(1.5), Inches(6.0), Inches(1.6)),
        ("🎙️", "Transcription & Summary",
         "Voice stories (Katha) → text transcript + AI-generated summary. Makes oral history searchable across generations.",
         NEELAM, Inches(0.5), Inches(3.3), Inches(6.0), Inches(1.6)),
        ("🗣️", "Digital Echo (Planned)",
         "RAG-based ancestor persona: bio + sample messages + relationships → AI speaks 'as your grandmother would have'. Explicit disclaimer.",
         PADMA, Inches(6.8), Inches(3.3), Inches(6.0), Inches(1.6)),
        ("🍳", "Recipe Generation",
         "Input ingredients → AI generates traditional Indian cooking steps, preserving family recipe knowledge even from partial memories.",
         PANNA, Inches(0.5), Inches(5.1), Inches(6.0), Inches(1.6)),
        ("🏷️", "Smart Tagging",
         "AI suggests culturally-relevant tags for content. Sentiment/emotion analysis on stories and messages.",
         SINDOOR, Inches(6.8), Inches(5.1), Inches(6.0), Inches(1.6)),
    ]

    for icon, title, desc, accent, x, y, w, h in features:
        _add_icon_card(s, x, y, w, h, icon, title, desc, accent)

    _add_slide_number(s, 10, TOTAL_SLIDES)


def slide_11_design_system(prs, blank):
    """Design System — Digital Sanskriti."""
    s = prs.slides.add_slide(blank)
    _add_bg(s, DARK_BG)
    _add_header_band(s, "Design System: Digital Sanskriti", "Every pixel inspired by Indian heritage", SUVARNA)

    # Color swatches
    _add_text(s, Inches(0.7), Inches(1.4), Inches(5), Inches(0.4), "Sacred Color Palette", 16, SUVARNA, True)

    colors = [
        ("Suvarna", "Temple Gold", SUVARNA),
        ("Sindoor", "Vermilion", SINDOOR),
        ("Padma", "Lotus Pink", PADMA),
        ("Haldi", "Turmeric", HALDI),
        ("Neelam", "Sapphire", NEELAM),
        ("Panna", "Emerald", PANNA),
        ("Khadi", "Aged Paper", KHADI),
        ("Masi", "Charcoal", RGBColor(60, 58, 55)),
    ]
    for i, (name, eng, color) in enumerate(colors):
        x = Inches(0.7 + (i % 4) * 1.5)
        y = Inches(1.9 + (i // 4) * 1.3)
        _add_shape(s, x, y, Inches(1.3), Inches(0.7), color)
        _add_text(s, x, y + Inches(0.75), Inches(1.3), Inches(0.25), name, 10, RGBColor(180, 175, 165), True, PP_ALIGN.CENTER)
        _add_text(s, x, y + Inches(0.95), Inches(1.3), Inches(0.2), eng, 9, MUTED, align=PP_ALIGN.CENTER)

    # Typography
    _add_text(s, Inches(0.7), Inches(4.8), Inches(5), Inches(0.4), "Typography", 16, SUVARNA, True)
    _add_bullets(s, Inches(0.7), Inches(5.25), Inches(5.5), Inches(2.0), [
        "Major Third scale (1.25 ratio): 10px → 60px",
        "Serif fonts: Georgia (iOS), Noto Serif (Web)",
        "Full Devanagari script support",
        "Handwritten style for personal notes",
    ], 13, RGBColor(200, 195, 185), Pt(4))

    # Motion system
    _add_text(s, Inches(6.8), Inches(1.4), Inches(6), Inches(0.4), "\"Digital Silk\" Motion System", 16, SUVARNA, True)

    motions = [
        ("silk", "Default transitions — smooth, luxurious"),
        ("water", "Flowing list animations"),
        ("pageFlip", "Card/page transitions"),
        ("heartbeat", "Pulse effects on selection"),
        ("dramatic", "Emphasis animations"),
        ("bounce", "Playful micro-interactions"),
        ("enter / exit", "Mount/unmount transitions"),
    ]
    for i, (name, desc) in enumerate(motions):
        _add_shape(s, Inches(6.8), Inches(1.9 + i * 0.55), Inches(5.8), Inches(0.45), RGBColor(55, 53, 50))
        _add_text(s, Inches(7.0), Inches(1.95 + i * 0.55), Inches(1.4), Inches(0.35), name, 12, SUVARNA, True)
        _add_text(s, Inches(8.4), Inches(1.95 + i * 0.55), Inches(4.0), Inches(0.35), desc, 11, RGBColor(170, 165, 155))

    # Spacing
    _add_text(s, Inches(6.8), Inches(5.8), Inches(6), Inches(0.4), "Spacing & Layout", 16, SUVARNA, True)
    _add_bullets(s, Inches(6.8), Inches(6.25), Inches(5.5), Inches(1.0), [
        "Golden Ratio-inspired from 4px base unit",
        "Gesture: pinch-zoom + pan + double-tap + haptics",
    ], 13, RGBColor(200, 195, 185), Pt(4))

    _add_slide_number(s, 11, TOTAL_SLIDES)


def slide_12_market(prs, blank):
    """Market Opportunity."""
    s = prs.slides.add_slide(blank)
    _add_bg(s, KHADI)
    _add_header_band(s, "Market Opportunity", "A massive, growing, underserved market", PANNA)

    # TAM/SAM/SOM
    _add_kpi_pill(s, Inches(0.5), Inches(1.5), "$8.2B by 2030", "Global Genealogy Market (CAGR ~10%)", PANNA, Inches(4.0))
    _add_kpi_pill(s, Inches(4.7), Inches(1.5), "300M+", "Indian Households (Primary Market)", SUVARNA_DK, Inches(4.0))
    _add_kpi_pill(s, Inches(8.9), Inches(1.5), "750M+", "Indian Smartphone Users", NEELAM, Inches(4.0))

    # Why Now
    _add_shape(s, Inches(0.5), Inches(2.9), Inches(6.0), Inches(4.3), CARD_BG, RGBColor(235, 228, 212))
    _add_text(s, Inches(0.7), Inches(3.0), Inches(5.6), Inches(0.4), "Why Now?", 18, PANNA, True)
    _add_bullets(s, Inches(0.7), Inches(3.5), Inches(5.6), Inches(3.5), [
        "📱  India crossed 800M smartphone users — mobile-first era",
        "🏠  Nuclear family trend → heritage loss accelerating",
        "😢  Post-COVID nostalgia driving cultural preservation interest",
        "🤖  AI (Gemini) makes photo analysis + persona chat viable at low cost",
        "💳  UPI + digital payments enable micro-subscriptions in India",
        "🌍  32M+ Indian diaspora seeking cultural roots",
        "📵  No competitor addresses Indian families specifically",
    ], 13, INK, Pt(5))

    # Target Segments
    _add_shape(s, Inches(6.8), Inches(2.9), Inches(6.0), Inches(4.3), CARD_BG, RGBColor(235, 228, 212))
    _add_text(s, Inches(7.0), Inches(3.0), Inches(5.6), Inches(0.4), "Target Segments", 18, SUVARNA_DK, True)

    segments = [
        ("🎯 Primary", "Indian families (urban + semi-urban) — 300M+ households"),
        ("🌍 Secondary", "Indian diaspora — 32M+ across 150 countries"),
        ("👴 Early Adopters", "Tech-savvy millennials who value family heritage"),
        ("🏛️ Institutional", "Cultural organizations, heritage foundations"),
    ]
    for i, (label, desc) in enumerate(segments):
        _add_text(s, Inches(7.0), Inches(3.5 + i * 0.9), Inches(5.6), Inches(0.3), label, 14, INK, True)
        _add_text(s, Inches(7.0), Inches(3.8 + i * 0.9), Inches(5.6), Inches(0.5), desc, 12, MUTED)

    _add_slide_number(s, 12, TOTAL_SLIDES)


def slide_13_competition(prs, blank):
    """Competitive Landscape."""
    s = prs.slides.add_slide(blank)
    _add_bg(s, KHADI)
    _add_header_band(s, "Competitive Landscape", "No direct competitor in the Indian heritage space", SINDOOR)

    # Table header
    headers = [("Feature", 0.20), ("Ancestry.com", 0.16), ("MyHeritage", 0.16), ("FamilySearch", 0.16), ("WhatsApp", 0.16), ("Vansh 🪷", 0.16)]
    _add_table_row_shapes(s, Inches(0.3), Inches(1.4), Inches(12.7), headers, Inches(0.4), SUVARNA_DK, WHITE, True, 12)

    rows = [
        ("Family Tree", "✅", "✅", "✅", "❌", "✅"),
        ("Indian Kinship Terms", "❌", "❌", "❌", "❌", "✅ 40+ terms"),
        ("Voice Stories", "❌", "❌", "❌", "❌ (audio msg)", "✅ + Transcript"),
        ("Photo Memories", "✅ (basic)", "✅", "✅ (basic)", "✅ (no context)", "✅ + AI Tags"),
        ("Traditions/Recipes", "❌", "❌", "❌", "❌", "✅"),
        ("Time-Locked Messages", "❌", "❌", "❌", "❌", "✅ (4 triggers)"),
        ("AI Features", "DNA matching", "Photo enhance", "❌", "❌", "✅ (5 AI tools)"),
        ("Offline-First", "❌", "❌", "Partial", "✅", "✅ (sync queue)"),
        ("Cultural Design", "Western", "Western", "Western", "Generic", "Digital Sanskriti"),
        ("Pricing", "$200+/yr", "$150+/yr", "Free (limited)", "Free", "Freemium ₹199/mo"),
        ("Mobile-First", "❌ (web)", "Partial", "❌ (web)", "✅", "✅ (native)"),
    ]

    for i, (feature, *vals) in enumerate(rows):
        bg = CARD_BG if i % 2 == 0 else KHADI
        cols = [(feature, 0.20)] + [(v, 0.16) for v in vals]
        _add_table_row_shapes(s, Inches(0.3), Inches(1.8 + i * 0.35), Inches(12.7), cols, Inches(0.35), bg, INK, False, 10)

    # Highlight vansh column
    _add_shape(s, Inches(10.98), Inches(1.4), Inches(2.02), Inches(0.05), PANNA, radius=False)

    _add_text(s, Inches(0.5), Inches(6.0), Inches(12), Inches(0.6),
              "Key Insight: No app combines family tree + voice stories + traditions + time-locked messages + AI — in an Indian cultural context.",
              14, SINDOOR, True)

    _add_slide_number(s, 13, TOTAL_SLIDES)


def slide_14_business_model(prs, blank):
    """Business Model."""
    s = prs.slides.add_slide(blank)
    _add_bg(s, KHADI)
    _add_header_band(s, "Business Model & Revenue", "Freemium with cultural value lock-in", SUVARNA)

    # Pricing tiers
    tiers = [
        ("Free", "₹0", ["1 family", "50 memories", "10 kathas", "Basic tree", "Offline mode"], MUTED),
        ("Legacy", "₹199/mo (~$2.50)", ["Unlimited memories", "AI features", "5GB storage", "Firebase sync", "Priority support"], SUVARNA_DK),
        ("Heritage", "₹499/mo (~$6)", ["Multiple families", "Digital Echo AI", "50GB storage", "Advanced analytics", "Custom themes"], PANNA),
    ]

    for i, (name, price, features, color) in enumerate(tiers):
        x = Inches(0.5 + i * 4.3)
        _add_shape(s, x, Inches(1.5), Inches(3.9), Inches(4.0), CARD_BG, RGBColor(235, 228, 212))
        _add_shape(s, x, Inches(1.5), Inches(3.9), Inches(0.85), color)
        _add_text(s, x + Inches(0.2), Inches(1.55), Inches(3.5), Inches(0.35), name, 20, WHITE, True, PP_ALIGN.CENTER)
        _add_text(s, x + Inches(0.2), Inches(1.9), Inches(3.5), Inches(0.35), price, 16, WHITE, False, PP_ALIGN.CENTER)
        _add_bullets(s, x + Inches(0.3), Inches(2.5), Inches(3.3), Inches(2.8), [f"✓  {f}" for f in features], 13, INK, Pt(6))

    # Key metrics
    _add_shape(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.5), CARD_BG, RGBColor(235, 228, 212))
    _add_text(s, Inches(0.7), Inches(5.8), Inches(11.9), Inches(0.4), "Key Metrics to Track", 16, SUVARNA_DK, True)

    metrics = [
        ("Families Created", "North Star metric"),
        ("Members/Family", "Virality indicator"),
        ("Content Uploads", "Engagement depth"),
        ("DAU/MAU", "Stickiness ratio"),
        ("Vasiyats Created", "Emotional lock-in"),
        ("Share Events", "Growth loop"),
    ]
    for i, (metric, desc) in enumerate(metrics):
        x = Inches(0.7 + (i % 6) * 2.05)
        _add_text(s, x, Inches(6.2), Inches(1.9), Inches(0.25), metric, 11, INK, True)
        _add_text(s, x, Inches(6.45), Inches(1.9), Inches(0.25), desc, 9, MUTED)

    _add_slide_number(s, 14, TOTAL_SLIDES)


def slide_15_closing(prs, blank):
    """Closing / Call to Action."""
    s = prs.slides.add_slide(blank)
    _add_bg(s, DARK_BG)

    _add_shape(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), SUVARNA, radius=False)

    _add_multi_text(s, Inches(1.5), Inches(1.0), Inches(10.5), Inches(2.0), [
        ("🪷", 50, SUVARNA, False),
        ("", 10, WHITE, False),
        ("Vansh — Every Family Has a Story Worth Preserving", 32, WHITE, True),
    ], PP_ALIGN.CENTER)

    # What's built
    _add_shape(s, Inches(1.0), Inches(3.2), Inches(5.5), Inches(3.0), RGBColor(55, 53, 50))
    _add_text(s, Inches(1.2), Inches(3.3), Inches(5.1), Inches(0.4), "✅  What's Built (POC)", 16, PANNA, True)
    _add_bullets(s, Inches(1.2), Inches(3.8), Inches(5.1), Inches(2.3), [
        "Full working mobile app (iOS/Android/Web)",
        "Complete REST API with 30+ endpoints",
        "Interactive family tree with custom algorithm",
        "Voice stories + photo memories + traditions",
        "Time-locked wisdom vault with encryption",
        "AI image analysis & era detection",
        "Offline-first with sync queue",
        "Indian kinship system (40+ terms)",
    ], 12, RGBColor(200, 195, 185), Pt(3))

    # Next steps
    _add_shape(s, Inches(6.8), Inches(3.2), Inches(5.5), Inches(3.0), RGBColor(55, 53, 50))
    _add_text(s, Inches(7.0), Inches(3.3), Inches(5.1), Inches(0.4), "🚀  Next Steps", 16, SUVARNA, True)
    _add_bullets(s, Inches(7.0), Inches(3.8), Inches(5.1), Inches(2.3), [
        "Deploy to cloud (AWS/GCP)",
        "App Store + Play Store submission",
        "Digital Echo AI persona launch",
        "DNA integration partnership",
        "Beta with 50 families",
        "Iterative UX improvements",
        "Marketing: heritage influencers",
        "Monetization: launch freemium tiers",
    ], 12, RGBColor(200, 195, 185), Pt(3))

    # Bottom
    _add_shape(s, Inches(0), Inches(6.8), SLIDE_W, Inches(0.7), SUVARNA_DK, radius=False)
    _add_text(s, Inches(1), Inches(6.88), Inches(6), Inches(0.4),
              "Thank You  •  Questions?", 18, WHITE, True)
    _add_text(s, Inches(7), Inches(6.88), Inches(5.5), Inches(0.4),
              "vansh-app  •  MIT License  •  Feb 2026", 14, WHITE, align=PP_ALIGN.RIGHT)

    _add_slide_number(s, 15, TOTAL_SLIDES)


# ═════════════════════════════════════════════════════════════════
# BUILD
# ═════════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    _set_widescreen(prs)
    blank = prs.slide_layouts[6]  # blank layout

    slide_01_title(prs, blank)
    slide_02_problem(prs, blank)
    slide_03_solution(prs, blank)
    slide_04_demo_home(prs, blank)
    slide_05_demo_tree(prs, blank)
    slide_06_architecture(prs, blank)
    slide_07_tech_deep_dive(prs, blank)
    slide_08_algorithms(prs, blank)
    slide_09_security(prs, blank)
    slide_10_ai_features(prs, blank)
    slide_11_design_system(prs, blank)
    slide_12_market(prs, blank)
    slide_13_competition(prs, blank)
    slide_14_business_model(prs, blank)
    slide_15_closing(prs, blank)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_FILE))
    print(f"\n🪷  Pitch deck saved to: {OUT_FILE}")
    print(f"    Slides: {TOTAL_SLIDES}")
    print(f"\n    Tip: Add screenshots to docs/ppt/images/ to replace placeholders!")


if __name__ == "__main__":
    build()
